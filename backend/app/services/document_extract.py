"""Extract text from uploaded documents (PDF, Office, images)."""

from __future__ import annotations

import base64
import csv
import io
import re
from typing import Any

import httpx
import structlog

from app.config import settings

log = structlog.get_logger("teambrain.document_extract")

DOC_TYPES = ("document", "field_report", "meeting_notes", "voice_note")

_EXT_MAP = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".csv": "csv",
    ".pptx": "pptx",
    ".txt": "txt",
    ".md": "txt",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".webp": "image",
    ".gif": "image",
    ".m4a": "audio",
    ".mp3": "audio",
    ".mpeg": "audio",
    ".ogg": "audio",
    ".wav": "audio",
    ".webm": "audio",
    ".aac": "audio",
}


def detect_format(filename: str, content_type: str | None) -> str:
    lower = (filename or "").lower()
    for ext, fmt in _EXT_MAP.items():
        if lower.endswith(ext):
            return fmt
    if content_type:
        ct = content_type.lower()
        if "pdf" in ct:
            return "pdf"
        if "word" in ct or "docx" in ct:
            return "docx"
        if "sheet" in ct or "excel" in ct:
            return "xlsx"
        if "presentation" in ct or "powerpoint" in ct:
            return "pptx"
        if "text" in ct:
            return "txt"
        if "image" in ct:
            return "image"
        if ct.startswith("audio/") or ct in ("application/ogg", "video/mp4"):
            return "audio"
    return "unknown"


def classify_doc_type(filename: str, text: str) -> str:
    """Auto-classify upload when possible."""
    lower = (filename or "").lower()
    if any(k in lower for k in ("rapport", "terrain", "mission", "field")):
        return "field_report"
    if any(k in lower for k in ("reunion", "meeting", "compte-rendu", "cr-")):
        return "meeting_notes"
    if any(k in lower for k in ("voice", "audio", "vocal")):
        return "voice_note"
    sample = (text or "")[:500].lower()
    if any(k in sample for k in ("mission terrain", "rapport terrain", "gps", "coordonnées")):
        return "field_report"
    if any(k in sample for k in ("décision", "réunion", "participants", "ordre du jour")):
        return "meeting_notes"
    return "document"


async def extract_text_from_bytes(
    content: bytes,
    filename: str,
    content_type: str | None = None,
) -> str:
    fmt = detect_format(filename, content_type)
    try:
        if fmt == "pdf":
            return _extract_pdf(content)
        if fmt == "docx":
            return _extract_docx(content)
        if fmt in ("xlsx", "csv"):
            return _extract_spreadsheet(content, fmt)
        if fmt == "pptx":
            return _extract_pptx(content)
        if fmt == "txt":
            return content.decode("utf-8", errors="replace")
        if fmt == "image":
            return await _extract_image_ocr(content, content_type or "image/jpeg")
    except Exception as exc:
        log.warning("extract_failed", filename=filename, fmt=fmt, error=str(exc))
    return ""


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts).strip()


def _extract_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text).strip()


def _extract_spreadsheet(content: bytes, fmt: str) -> str:
    if fmt == "csv":
        text = content.decode("utf-8", errors="replace")
        rows = list(csv.reader(io.StringIO(text)))
        return "\n".join(" | ".join(row) for row in rows[:200])
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"## {sheet.title}")
        for i, row in enumerate(sheet.iter_rows(values_only=True)):
            if i > 200:
                break
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        if slide_text:
            parts.append(f"Slide {i}:\n" + "\n".join(slide_text))
    return "\n\n".join(parts).strip()


async def _extract_image_ocr(content: bytes, content_type: str) -> str:
    if not settings.gemini_api_key:
        return ""
    b64 = base64.b64encode(content).decode("ascii")
    mime = content_type if content_type.startswith("image/") else "image/jpeg"
    url = (
        "https://generativelanguage.googleapis.com/v1beta/models/"
        f"gemini-2.5-flash:generateContent?key={settings.gemini_api_key}"
    )
    body: dict[str, Any] = {
        "contents": [
            {
                "parts": [
                    {"text": "Extrais tout le texte visible de cette image. Réponds uniquement avec le texte extrait."},
                    {"inline_data": {"mime_type": mime, "data": b64}},
                ]
            }
        ],
        "generationConfig": {"temperature": 0.1},
    }
    async with httpx.AsyncClient(timeout=90) as client:
        r = await client.post(url, json=body)
        if r.status_code != 200:
            return ""
        data = r.json()
        return data["candidates"][0]["content"]["parts"][0]["text"].strip()


def parse_tags_from_text(text: str) -> list[str]:
    words = re.findall(r"#(\w+)", text)
    return list(dict.fromkeys(words))[:10]
